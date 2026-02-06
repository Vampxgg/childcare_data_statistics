# -*- coding: utf-8 -*-
"""
节点ID: 1754713960183
节点标题: 文档提取
节点描述: 
节点类型: code
"""

import requests
import io
import os
from urllib.parse import urlparse
import time # 引入time模块，可以在批处理之间添加延时
 
# 各文件类型的处理库
import docx
import pandas as pd
from pptx import Presentation
import pdfplumber
 
def extract_content_from_url(url: str) -> str:
    """
    根据给定的URL，下载文件并提取其文本内容。
    支持的格式: .txt, .csv, .json, .md, .docx, .xlsx, .pptx, .pdf
 
    Args:
        url (str): 文件的公开URL链接。
 
    Returns:
        str: 提取出的文本内容，如果失败或不支持则返回错误信息。
    """
    print(f"[*] 正在处理URL: {url}")
 
    try:
        # 1. 下载文件
        # 添加headers模拟浏览器，可以提高成功率
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        response = requests.get(url, timeout=30, headers=headers)
        response.raise_for_status() # 检查请求是否成功 (非2xx状态码会引发HTTPError)
        print("[+] 文件下载成功")
 
        file_in_memory = io.BytesIO(response.content)
 
        # 2. 识别文件类型
        path = urlparse(url).path
        ext = os.path.splitext(path)[1].lower()
        
        # 如果URL没有后缀，尝试从响应头中识别
        if not ext:
            content_type = response.headers.get('Content-Type', '')
            print(f"[*] URL无后缀，检查Content-Type: {content_type}")
            if 'pdf' in content_type: ext = '.pdf'
            elif 'msword' in content_type: ext = '.doc' # 注意，这个库不支持.doc
            elif 'vnd.openxmlformats-officedocument.wordprocessingml.document' in content_type: ext = '.docx'
            elif 'vnd.openxmlformats-officedocument.spreadsheetml.sheet' in content_type: ext = '.xlsx'
            elif 'vnd.openxmlformats-officedocument.presentationml.presentation' in content_type: ext = '.pptx'
 
        print(f"[*] 文件类型识别为: {ext}")
        
        # 3. 根据不同类型进行内容提取
        content = ""
        if ext in ['.txt', '.csv', '.json', '.md', '.html', '.xml']:
            content = response.content.decode('utf-8', errors='ignore')
 
        elif ext == '.docx':
            document = docx.Document(file_in_memory)
            paragraphs = [p.text for p in document.paragraphs]
            content = "\n".join(paragraphs)
 
        elif ext == '.xlsx':
            xls = pd.ExcelFile(file_in_memory)
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                content += f"--- 工作表: {sheet_name} ---\n{df.to_string()}\n\n"
 
        elif ext == '.pptx':
            prs = Presentation(file_in_memory)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
 
        elif ext == '.pdf':
            with pdfplumber.open(file_in_memory) as pdf:
                all_text = [page.extract_text() for page in pdf.pages if page.extract_text()]
                content = "\n".join(all_text)
                        
        else:
            return f"错误：不支持或无法识别的文件类型 '{ext}'。"
 
        print("[+] 内容提取完成！")
        return content.strip()
 
    except requests.exceptions.HTTPError as e:
        return f"错误：HTTP请求失败 - {e.response.status_code} {e.response.reason} for URL: {url}"
    except requests.exceptions.RequestException as e:
        return f"错误：网络请求失败 - {e}"
    except Exception as e:
        return f"错误：处理文件时发生未知错误 - {e}"
 
# --- 您要求的批处理函数 ---
 
def main(files: list) -> dict:
    """
    批量处理文件对象列表，提取每个文件URL指向的内容。
 
    Args:
        files (list): 对象列表，每个对象都应有 .url 属性。
                      例如: [File(url='http://...'), File(url='http://...')]
 
    Returns:
        dict: 包含所有处理结果的字典，格式如下:
              {
                  "result": [
                      {"url": "url1", "content": "extracted_content1"},
                      {"url": "url2", "content": "error_message2"},
                      ...
                  ]
              }
    """
    results = []
    arry_str=[]
    total_files = len(files)
    print(f"--- 开始批量处理 {total_files} 个文件 ---")
 
    for i, file_obj in enumerate(files):
        print(f"\n--- 正在处理第 {i + 1}/{total_files} 个文件 ---")
        
        # 健壮性检查：确保对象有url属性且不为空
        # if not hasattr(file_obj, 'url') or not getattr(file_obj, 'url'):
        #     print("[!] 警告：输入对象格式不正确或URL为空，已跳过。")
        #     result_item = {
        #         "url": getattr(file_obj, 'url', 'URL缺失'),
        #         "content": "错误：输入对象格式不正确或URL为空。"
        #     }
        #     results.append(result_item)
        #     continue
        
        url = file_obj['url']
        extracted_content = extract_content_from_url(url)
        
        result_item = {
            "url": url,
            "content": extracted_content
        }
        results.append(result_item)
        arry_str.append(extracted_content)
        # 礼貌性地暂停一下，避免对服务器造成冲击
        # time.sleep(1) 
 
    # print("\n--- 所有文件处理完毕 ---")
    # 返回您指定的最终格式
    return {"result": results,
            "arry_str":arry_str}
