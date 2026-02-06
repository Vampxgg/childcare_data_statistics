# -*- coding: utf-8 -*-
"""
节点ID: 17617451066163
节点标题: 文档提取 (1)
节点描述: 
节点类型: code
"""

import requests
import io
import os
from urllib.parse import urlparse
import time
 
# --- 各文件类型的处理库 (根据您提供的列表) ---
import docx
import pandas as pd
from pptx import Presentation
import pdfplumber
 
def extract_content_from_url(url: str) -> str:
    """
    根据给定的URL，下载文件并提取其文本内容。
    使用您指定的原始库。
    支持的格式: .txt, .csv, .json, .md, .docx, .xlsx, .pptx, .pdf
 
    Args:
        url (str): 文件的公开URL链接。
 
    Returns:
        str: 提取出的文本内容，如果失败或不支持则返回错误信息。
    """
    print(f"[*] 正在处理URL: {url}")
    try:
        headers = { 'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36' }
        response = requests.get(url, timeout=60, headers=headers)
        response.raise_for_status()
        print("[+] 文件下载成功")
        file_in_memory = io.BytesIO(response.content)
 
        path = urlparse(url).path
        ext = os.path.splitext(path)[1].lower()
        
        if not ext:
            content_type = response.headers.get('Content-Type', '')
            print(f"[*] URL无后缀，检查Content-Type: {content_type}")
            if 'pdf' in content_type: ext = '.pdf'
            elif 'msword' in content_type: ext = '.doc' # 注意，这个库不支持.doc
            elif 'vnd.openxmlformats-officedocument.wordprocessingml.document' in content_type: ext = '.docx'
            elif 'vnd.openxmlformats-officedocument.spreadsheetml.sheet' in content_type: ext = '.xlsx'
            elif 'vnd.openxmlformats-officedocument.presentationml.presentation' in content_type: ext = '.pptx'
 
        print(f"[*] 文件类型识别为: {ext}")
        
        content = ""
        if ext in ['.txt', '.csv', '.json', '.md', '.html', '.xml']:
            content = response.content.decode('utf-8', errors='ignore')
 
        elif ext == '.docx':
            # 注意：此方法可能因处理含SVG等复杂元素的文档而失败
            document = docx.Document(file_in_memory)
            paragraphs = [p.text for p in document.paragraphs]
            content = "\n".join(paragraphs)
 
        elif ext == '.xlsx':
            # 需要安装 openpyxl
            xls = pd.ExcelFile(file_in_memory, engine='openpyxl')
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                content += f"--- 工作表: {sheet_name} ---\n{df.to_string()}\n\n"
 
        elif ext == '.pptx':
            prs = Presentation(file_in_memory)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
 
        elif ext == '.pdf':
            with pdfplumber.open(file_in_memory) as pdf:
                all_text = [page.extract_text() for page in pdf.pages if page.extract_text()]
                content = "\n".join(all_text)
                        
        else:
            return f"错误：不支持或无法识别的文件类型 '{ext}'。"
 
        print("[+] 内容提取完成！")
        return content.strip()
 
    except requests.exceptions.HTTPError as e:
        return f"错误：HTTP请求失败 - {e.response.status_code} {e.response.reason} for URL: {url}"
    except requests.exceptions.RequestException as e:
        return f"错误：网络请求失败 - {e}"
    except Exception as e:
        # 添加更详细的错误追踪信息，便于调试
        import traceback
        print(traceback.format_exc())
        return f"错误：处理文件时发生未知错误 - {e}"
 
# --- 核心修改：实现分类逻辑的 main 函数 ---
 
def main(files: list) -> dict:
    """
    批量处理文件字典列表，根据文件ID将提取的内容分类。
 
    Args:
        files (list): 字典列表，每个字典必须包含 'id' 和 'url' 键。
 
    Returns:
        dict: 包含分类后文本内容列表的字典，格式如下:
              {
                  "industrial_area_report": [...],
                  "talent_demand_analysis": [...],
                  "old_talent_development_plan": [...]
              }
    """
    # 1. 初始化用于存储最终结果的列表
    industrial_area_report = []
    talent_demand_analysis = []
    old_talent_development_plan = []
 
    total_files = len(files)
    print(f"--- 开始批量处理 {total_files} 个文件 ---")
 
    # 2. 遍历输入的每个文件信息字典
    for i, file_info in enumerate(files):
        print(f"\n--- 正在处理第 {i + 1}/{total_files} 个文件 ---")
        
        # 使用 .get() 安全地获取 'id' 和 'url'，防止因键不存在而报错
        file_id = file_info.get('id')
        url = file_info.get('url')
        # file_id = file_info['id']
        # url = file_info['url']
 
        # 健壮性检查：如果缺少关键信息，则跳过此文件
        if not file_id or not url:
            print(f"[!] 警告: 跳过一个条目，因为它缺少 'id' 或 'url'。条目: {file_info}")
            continue
 
        # 3. 调用函数提取文件内容
        extracted_content = extract_content_from_url(url)
        
        # 4. 根据 'id' 的前缀进行分类
        if file_id.startswith('1'):
            print(f"[*] 文件 ID '{file_id}' 匹配规则 '1', 内容归入 'industrial_area_report'")
            industrial_area_report.append(extracted_content)
        elif file_id.startswith('2'):
            print(f"[*] 文件 ID '{file_id}' 匹配规则 '2', 内容归入 'talent_demand_analysis'")
            talent_demand_analysis.append(extracted_content)
        elif file_id.startswith('3'):
            print(f"[*] 文件 ID '{file_id}' 匹配规则 '3', 内容归入 'old_talent_development_plan'")
            old_talent_development_plan.append(extracted_content)
        else:
            print(f"[!] 警告: 文件 ID '{file_id}' (URL: {url}) 没有匹配任何分类规则，已忽略。")
        # industrial_area_report.append(extracted_content)
        # talent_demand_analysis.append(extracted_content)
        # old_talent_development_plan.append(extracted_content)



        # 礼貌性地暂停，避免对服务器造成过大压力
        # time.sleep(0.5)
 
    print("\n--- 所有文件处理完毕 ---")
 
    # 5. 按要求的格式返回最终的字典
    return {
        "industrial_area_report": industrial_area_report,
        "talent_demand_analysis": talent_demand_analysis,
        "old_talent_development_plan": old_talent_development_plan
    }
